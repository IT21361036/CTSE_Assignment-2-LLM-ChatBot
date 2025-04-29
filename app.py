# This chatbot answers questions based on CTSE lecture notes using OpenRouter's LLM API.
import os
import openai
from PyPDF2 import PdfReader
from typing import List
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import ipywidgets as widgets
from IPython.display import display, Markdown, clear_output
import streamlit as st
from streamlit_chat import message
from pptx import Presentation  # For PowerPoint files
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables import RunnablePassthrough
from langchain_core.messages import HumanMessage, AIMessage
from langchain_core.chat_history import InMemoryChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory

# Configuration
class Config:
    OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "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")
    OPENROUTER_BASE_URL = "https://openrouter.ai/api/v1" # Replace with your actual API key
   # OPENROUTER_BASE_URL = "https://openrouter.ai/api/v1"
    MODEL_NAME = "mistralai/mistral-7b-instruct"  # Cost-effective model
    TEMPERATURE = 0.3  # Controls randomness of responses
    MAX_TOKENS = 1000  # Limit response length
    CONTEXT_TOKENS = 3000  # Max context to send to LLM
    CHUNK_SIZE = 500  # Size of text chunks for processing
# Initialize OpenAI client for OpenRouter
client = openai.OpenAI(
    base_url=Config.OPENROUTER_BASE_URL,
    api_key=Config.OPENROUTER_API_KEY,
)

# Document Processing Functions
def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from a PDF file"""
    text = ""
    try:
        with open(pdf_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""  # Handle None returns
    except Exception as e:
        print(f"Error reading PDF {pdf_path}: {str(e)}")
    return text
#new
def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from a PowerPoint file"""
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {pptx_path}: {str(e)}")
    return text

def chunk_text(text: str, chunk_size: int = Config.CHUNK_SIZE) -> List[str]:
    """Split text into manageable chunks"""
    words = text.split()
    chunks = [' '.join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
    return chunks

def create_vector_index(text_chunks: List[str]):
    """Create TF-IDF vector index for semantic search"""
    vectorizer = TfidfVectorizer(stop_words='english')
    tfidf_matrix = vectorizer.fit_transform(text_chunks)
    return vectorizer, tfidf_matrix

# Knowledge Base Class
class LectureNotesKB:
    def __init__(self):
        self.notes = {}
        self.vectorizers = {}
        self.tfidf_matrices = {}
        self.chunks = {}
    
    def add_lecture(self, lecture_name: str, file_path: str):
        """Add a lecture to the knowledge base from a PDF or PPTX"""
        text = ""
        if file_path.lower().endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        elif file_path.lower().endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        else:
            print(f"Unsupported file type for {file_path}")
            return
        
        if not text:
            print(f"Warning: No text extracted from {file_path}")
            return
            
        chunks = chunk_text(text)
        vectorizer, tfidf_matrix = create_vector_index(chunks)
        
        self.notes[lecture_name] = text
        self.chunks[lecture_name] = chunks
        self.vectorizers[lecture_name] = vectorizer
        self.tfidf_matrices[lecture_name] = tfidf_matrix
        print(f"Loaded lecture: {lecture_name} ({len(chunks)} chunks)")
     
    def get_relevant_chunks(self, lecture_name: str, query: str, top_k: int = 3) -> List[str]:
        """Retrieve most relevant text chunks for a query"""
        if lecture_name not in self.notes:
            print(f"Lecture {lecture_name} not found in knowledge base")
            return []
            
        vectorizer = self.vectorizers[lecture_name]
        tfidf_matrix = self.tfidf_matrices[lecture_name]
        
        query_vec = vectorizer.transform([query])
        similarities = cosine_similarity(query_vec, tfidf_matrix).flatten()
        top_indices = similarities.argsort()[-top_k:][::-1]
        
        return [self.chunks[lecture_name][i] for i in top_indices]
# Chatbot Class with LangChain
class LectureChatbot:
    def __init__(self):
        self.knowledge_base = LectureNotesKB()
        # Initialize LangChain LLM
        self.llm = ChatOpenAI(
            model=Config.MODEL_NAME,
            openai_api_key=Config.OPENROUTER_API_KEY,
            openai_api_base=Config.OPENROUTER_BASE_URL,
            temperature=Config.TEMPERATURE,
            max_tokens=Config.MAX_TOKENS
        )
        # Define prompt template
        self.prompt = ChatPromptTemplate.from_messages([
            ("system", """You are an expert teaching assistant for the CTSE course. 
Answer the student's question based strictly on the provided lecture notes context.
If the answer isn't in the notes, say "I don't have that information in my lecture notes."

Lecture Context:
{context}"""),
            MessagesPlaceholder(variable_name="history"),
            ("human", "{input}")
        ])
        # Initialize chat history store
        self.chat_histories = {}
        # Create runnable chain
        self.chain = (
            RunnablePassthrough.assign(
                context=lambda x: "\n\n".join(
                    self.knowledge_base.get_relevant_chunks(x["lecture_name"], x["input"])
                )[:Config.CONTEXT_TOKENS]
            )
            | self.prompt
            | self.llm
        )
        # Wrap with message history
        self.runnable = RunnableWithMessageHistory(
            self.chain,
            lambda session_id: self.get_chat_history(session_id),
            input_messages_key="input",
            history_messages_key="history"
        )

    def get_chat_history(self, session_id: str) -> InMemoryChatMessageHistory:
        """Get or create chat history for a session"""
        if session_id not in self.chat_histories:
            self.chat_histories[session_id] = InMemoryChatMessageHistory()
        return self.chat_histories[session_id]

    def load_lecture(self, lecture_name: str, file_path: str):
        """Load lecture notes (PDF or PPTX) into the knowledge base"""
        self.knowledge_base.add_lecture(lecture_name, file_path)

    def generate_response(self, lecture_name: str, question: str, session_id: str) -> str:
        """Generate a response to the user's question"""
        try:
            response = self.runnable.invoke(
                {"input": question, "lecture_name": lecture_name},
                config={"configurable": {"session_id": session_id}}
            )
            return response.content
        except Exception as e:
            return f"Error generating response: {str(e)}"

    def clear_history(self, session_id: str):
        """Clear chat history for a session"""
        if session_id in self.chat_histories:
            self.chat_histories[session_id] = InMemoryChatMessageHistory()

# Streamlit Web Interface
def main():
    st.set_page_config(page_title="CTSE Lecture Chatbot", page_icon="📚", layout="wide")
    st.title("CTSE Lecture Chatbot")
    
    chatbot = LectureChatbot()
    LECTURE_NOTES_DIR = "ctse_lecture_notes"
    os.makedirs(LECTURE_NOTES_DIR, exist_ok=True)
    
    # File upload
    uploaded_file = st.file_uploader("Upload a PDF or PPTX lecture", type=["pdf", "pptx"])
    if uploaded_file:
        file_path = os.path.join(LECTURE_NOTES_DIR, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        chatbot.load_lecture(uploaded_file.name.split('.')[0].title(), file_path)
        st.success(f"Uploaded and loaded: {uploaded_file.name}")
    
    # Scan directory for PDF and PPTX files
    available_lectures = {}
    try:
        for file in os.listdir(LECTURE_NOTES_DIR):
            if file.lower().endswith(('.pdf', '.pptx')):
                lecture_name = os.path.splitext(file)[0].replace('_', ' ').title()
                available_lectures[lecture_name] = os.path.join(LECTURE_NOTES_DIR, file)
                
        if not available_lectures:
            st.error(f"No PDF or PPTX files found in {LECTURE_NOTES_DIR}. Please upload or add lecture files.")
            return
            
        # Load available lectures
        for name, path in available_lectures.items():
            if name not in chatbot.knowledge_base.notes:
                chatbot.load_lecture(name, path)
        
        # Lecture selection
        lecture_name = st.selectbox("Select Lecture:", list(available_lectures.keys()), key="lecture_select")
        
        # Initialize session state for chat history
        if "messages" not in st.session_state:
            st.session_state.messages = []
        
        # Session ID for chat history
        session_id = f"{lecture_name}_session"
        
        # Clear chat button
        if st.button("Clear Chat"):
            st.session_state.messages = []
            chatbot.clear_history(session_id)
            st.rerun()
        
        # Display chat history
        for idx, msg in enumerate(st.session_state.messages):
            if msg["lecture"] == lecture_name:
                message(msg["content"], is_user=(msg["role"] == "user"), key=f"msg_{idx}")
        
        # Input form for question
        with st.form(key="question_form", clear_on_submit=True):
            question = st.text_input("Ask a question about the lecture...", key="question_input")
            submit_button = st.form_submit_button("Ask")
            
            if submit_button and question.strip():
                # Add user message
                st.session_state.messages.append({"role": "user", "content": question, "lecture": lecture_name})
                # Generate response
                answer = chatbot.generate_response(lecture_name, question, session_id)
                # Add assistant response
                st.session_state.messages.append({"role": "assistant", "content": answer, "lecture": lecture_name})
                # Rerun to update UI
                st.rerun()
                
    except Exception as e:
        st.error(f"Error loading lectures from {LECTURE_NOTES_DIR}: {str(e)}")
# Run the chatbot
if __name__ == "__main__":
    main()
